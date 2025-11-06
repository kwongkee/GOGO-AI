# -*- coding: utf-8 -*-
import os
import sys
import json
import redis
import hashlib
import psutil
from datetime import datetime, timedelta
from celery import Celery
import docx
import PyPDF2
from openpyxl import load_workbook
import xlrd
from pptx import Presentation
import traceback
import re
import numpy as np
from celery.exceptions import Reject
from transformers import AutoTokenizer, AutoModel
import torch
import logging
import threading
import tempfile
from PIL import Image
#import whisper
from transformers import WhisperProcessor, WhisperForConditionalGeneration
import torchaudio
import requests
from unstructured.partition.auto import partition
from moviepy.editor import VideoFileClip
from sentence_transformers import SentenceTransformer
from kafka import KafkaProducer

# 全局生产者
producer = None

import jieba
import warnings
warnings.filterwarnings("ignore")
jieba.setLogLevel(logging.CRITICAL)

def get_producer():
    global producer
    if producer is None:
        producer = KafkaProducer(
            bootstrap_servers=['39.108.11.214:9092'],
            value_serializer=lambda v: json.dumps(v, ensure_ascii=False).encode('utf-8'),
            acks=1,
            retries=3,
            batch_size=16384,
            linger_ms=10
        )
    return producer
    
# 配置模型路径
MODEL_CACHE_DIR = "/var/cache/huggingface/"
os.makedirs(MODEL_CACHE_DIR, exist_ok=True)

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("/var/log/multimodal_rag/vectorize.log")
    ]
)
logger = logging.getLogger(__name__)

# 关闭所有非必要日志（只保留 CRITICAL）
logging.getLogger().setLevel(logging.CRITICAL)

class ModelManager:
    """模型管理器，负责模型的加载、卸载和健康检查"""
    
    _instance = None
    _lock = threading.Lock()
    
    def __new__(cls):
        if cls._instance is None:
            with cls._lock:
                if cls._instance is None:
                    cls._instance = super().__new__(cls)
                    cls._instance._initialized = False
        return cls._instance
    
    def __init__(self):
        if not self._initialized:
            self.clip_model = None
            self.clip_processor = None
            self.whisper_model = None
            self.whisper_processor = None  # 添加这行
            self.bert_model = None
            self.bert_tokenizer = None
            self._load_lock = threading.Lock()
            self._models_loaded = False
            self._initialized = True
    
    def load_models(self):
        """加载所有需要的模型"""
        with self._load_lock:
            if self._models_loaded:
                logger.info("模型已加载，跳过重复加载")
                return True
                
            try:
                if self.clip_model is None:
                    logger.info("加载CLIP模型...")
                    local_clip_path = "/var/cache/huggingface/chinese-clip-vit-large-patch14"
                    if not os.path.exists(local_clip_path):
                        logger.info("本地路径不存在，正在从镜像下载...")
                        snapshot_download(repo_id="OFA-Sys/chinese-clip-vit-large-patch14", local_dir=local_clip_path)
                    from transformers import ChineseCLIPModel, ChineseCLIPProcessor
                    self.clip_model = ChineseCLIPModel.from_pretrained(
                        local_clip_path,
                        local_files_only=True,
                        trust_remote_code=True
                    )
                    self.clip_processor = ChineseCLIPProcessor.from_pretrained(
                        local_clip_path,
                        local_files_only=True,
                        trust_remote_code=True
                    )
                    logger.info("CLIP模型加载成功")
                
                if self.whisper_model is None:
                    logger.info("加载Whisper模型...")
                    local_whisper_path = "/var/cache/huggingface/whisper-medium"
                    
                    if not os.path.exists(local_whisper_path):
                        logger.warning("Whisper模型不存在，跳过音频处理功能")
                        self.whisper_model = None
                        self.whisper_processor = None  # 确保这里也设置
                    else:
                        try:
                            self.whisper_processor = WhisperProcessor.from_pretrained(
                                local_whisper_path,
                                local_files_only=True
                            )
                            self.whisper_model = WhisperForConditionalGeneration.from_pretrained(
                                local_whisper_path,
                                local_files_only=True
                            )
                            self.whisper_model.eval()
                            logger.info("Whisper模型加载成功")
                        except Exception as e:
                            logger.warning(f"Whisper模型加载失败: {e}，跳过音频处理")
                            self.whisper_model = None
                            self.whisper_processor = None  # 确保这里也设置
                
                if self.bert_model is None:
                    logger.info("加载BERT模型...")
                    model_path = "/var/cache/huggingface/text2vec-base-chinese"
                    if not os.path.exists(model_path):
                        logger.info("本地路径不存在，正在从镜像下载...")
                        snapshot_download(repo_id="shibing624/text2vec-base-chinese", local_dir=model_path)
                    self.bert_model = AutoModel.from_pretrained(model_path)
                    self.bert_tokenizer = AutoTokenizer.from_pretrained(model_path)
                
                self._models_loaded = True
                logger.info("所有模型首次加载完成")
                return True
                
            except Exception as e:
                logger.error(f"模型加载失败: {e}")
                self._models_loaded = False
                raise
    
    def unload_models(self):
        """卸载模型以释放内存"""
        with self._load_lock:
            self.clip_model = None
            self.whisper_model = None
            self.bert_model = None
            self.bert_tokenizer = None
            self._models_loaded = False  # 重置实例变量
            logger.info("模型已卸载")
            
    def health_check(self):
        """模型健康检查 - 修复属性访问"""
        try:
            return {
                'clip_model': self.clip_model is not None,
                'whisper_model': self.whisper_model is not None,
                'whisper_processor': hasattr(self, 'whisper_processor') and self.whisper_processor is not None,  # 安全访问
                'bert_model': self.bert_model is not None,
                'bert_tokenizer': self.bert_tokenizer is not None
            }
        except AttributeError as e:
            logger.warning(f"健康检查属性访问错误: {e}")
            # 返回所有False，触发重新加载
            return {
                'clip_model': False,
                'whisper_model': False,
                'whisper_processor': False,
                'bert_model': False,
                'bert_tokenizer': False
            }

    def ensure_nltk_data(self):
        """确保NLTK数据包存在"""
        try:
            import nltk
            import ssl
            
            # 解决SSL证书问题
            try:
                _create_unverified_https_context = ssl._create_unverified_context
            except AttributeError:
                pass
            else:
                ssl._create_default_https_context = _create_unverified_https_context
            
            # 检查并下载必要的数据包
            required_packages = ['punkt_tab', 'averaged_perceptron_tagger_eng', 'punkt']
            for package in required_packages:
                try:
                    nltk.data.find(f'tokenizers/{package}')
                except LookupError:
                    logger.info(f"下载NLTK数据包: {package}")
                    nltk.download(package, quiet=True)
                    
        except Exception as e:
            logger.warning(f"NLTK数据包检查失败: {e}")
    
class ConfigManager:
    """配置管理器"""
    
    def __init__(self, config_path='/app/config.ini'):
        self.config_path = config_path
        self.default_config = {
            'text_weight': 0.7,
            'video_weight': 0.6,
            'image_weight': 0.8,
            'multi_element_adjust': 0.1,
            'max_segment_length': 2048,
            'vector_dimension': 768,
            'max_file_size': 100 * 1024 * 1024,  # 100MB
            'supported_formats': {
                'documents': ['.pdf', '.pptx', '.docx', '.xls', '.xlsx'],
                'images': ['.jpg', '.png', '.jpeg', '.bmp', '.gif'],
                'audio': ['.mp3', '.wav', '.m4a'],
                'video': ['.mp4', '.avi', '.mov'],
                'text': ['.txt', '.csv', '.json']
            }
        }
        self.config = self._load_config()
    
    def _load_config(self):
        """加载配置文件"""
        config = self.default_config.copy()
        
        try:
            if os.path.exists(self.config_path):
                import configparser
                parser = configparser.ConfigParser()
                parser.read(self.config_path)
                
                # 读取模型配置
                if parser.has_section('Models'):
                    model_config = dict(parser.items('Models'))
                    # 可以在这里添加模型特定配置
                
                # 读取处理配置
                if parser.has_section('Processing'):
                    processing_config = dict(parser.items('Processing'))
                    config.update({
                        'text_weight': float(processing_config.get('text_weight', 0.7)),
                        'video_weight': float(processing_config.get('video_weight', 0.6)),
                        'image_weight': float(processing_config.get('image_weight', 0.8)),
                        'multi_element_adjust': float(processing_config.get('multi_element_adjust', 0.1)),
                        'max_segment_length': int(processing_config.get('max_segment_length', 2048)),
                        'max_file_size': int(processing_config.get('max_file_size', 100 * 1024 * 1024))
                    })
                
                # 读取系统配置
                if parser.has_section('System'):
                    system_config = dict(parser.items('System'))
                    # 可以在这里添加系统配置
                    
        except Exception as e:
            logger.warning(f"配置文件加载失败: {e}, 使用默认配置")
        
        return config
    
    def get(self, key, default=None):
        """获取配置值"""
        return self.config.get(key, default)
    
    def update(self, updates):
        """更新配置（不持久化）"""
        self.config.update(updates)

class VectorQualityValidator:
    """向量质量验证器"""
    
    @staticmethod
    def validate_vector(vector, original_content=""):
        """验证向量质量"""
        if not vector or len(vector) == 0:
            return False, "空向量"
        
        vector_np = np.array(vector)
        
        # 检查向量是否全为零
        if np.all(vector_np == 0):
            return False, "全零向量"
        
        # 检查向量范数
        norm = np.linalg.norm(vector_np)
        if norm < 1e-6:
            return False, "向量范数过小"
        
        # 检查向量是否包含NaN或Inf
        if np.any(np.isnan(vector_np)) or np.any(np.isinf(vector_np)):
            return False, "向量包含NaN或Inf"
        
        # 检查向量维度一致性
        if len(vector) != 768:  # 假设标准维度为768
            return False, f"向量维度不一致: {len(vector)} != 768"
        
        return True, "质量合格"

class FileProcessor:
    """文件处理器"""
    
    def __init__(self, config_manager):
        self.config = config_manager
        self.model_manager = ModelManager()
    
    def validate_file(self, file_path):
        """验证文件"""
        if not os.path.exists(file_path):
            return False, "文件不存在"
        
        if not os.path.isfile(file_path):
            return False, "不是文件"
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        if file_size > self.config.get('max_file_size'):
            return False, f"文件过大: {file_size} bytes"
        
        # 检查文件格式
        file_ext = os.path.splitext(file_path)[1].lower()
        supported = False
        for category, formats in self.config.get('supported_formats', {}).items():
            if file_ext in formats:
                supported = True
                break
        
        if not supported:
            return False, f"不支持的文件格式: {file_ext}"
        
        return True, "验证通过"
    
    def extract_content(self, file_path):
        """提取文件内容"""
        valid, message = self.validate_file(file_path)
        if not valid:
            return "", [], message
        
        file_ext = os.path.splitext(file_path)[1].lower()
        text_content = ""
        multimodal_elements = []
        
        try:
            if file_ext in self.config.get('supported_formats', {}).get('documents', []):
                text_content, multimodal_elements = self._process_document(file_path)
                message = "处理成功"
            elif file_ext in self.config.get('supported_formats', {}).get('images', []):
                multimodal_elements = self._process_image(file_path)
                message = "处理成功"
            elif file_ext in self.config.get('supported_formats', {}).get('audio', []):
                text_content = self._process_audio(file_path)
                message = "处理成功"
            elif file_ext in self.config.get('supported_formats', {}).get('video', []):
                text_content, multimodal_elements, _ = self._process_video(file_path)
                message = "处理成功"
            elif file_ext in self.config.get('supported_formats', {}).get('text', []):
                text_content = self._process_text(file_path)
                message = "处理成功"
            else:
                return "", [], f"未实现的处理方法: {file_ext}"
                
            return text_content, multimodal_elements, message
            
        except Exception as e:
            logger.error(f"文件处理失败: {file_path}, 错误: {e}")
            return "", [], f"处理失败: {str(e)}"
    
    def _process_document(self, file_path):
        """处理文档文件"""
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            text_content = ""
            multimodal_elements = []
            
            if file_ext == '.docx':
                # 使用python-docx直接处理
                doc = docx.Document(file_path)
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_content += paragraph.text + "\n"
                        
            elif file_ext == '.pdf':
                # 使用PyPDF2处理PDF
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text_content += page.extract_text() + "\n"
                        
            elif file_ext in ['.xls', '.xlsx']:
                # 处理Excel文件
                try:
                    wb = load_workbook(file_path)
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        for row in ws.iter_rows(values_only=True):
                            row_text = ' '.join([str(cell) for cell in row if cell])
                            if row_text.strip():
                                text_content += row_text + "\n"
                except:
                    # 备选方案：使用xlrd
                    workbook = xlrd.open_workbook(file_path)
                    for sheet in workbook.sheets():
                        for row in range(sheet.nrows):
                            row_text = ' '.join([str(sheet.cell_value(row, col)) for col in range(sheet.ncols)])
                            if row_text.strip():
                                text_content += row_text + "\n"
                                
            elif file_ext == '.pptx':
                # 处理PPT文件
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_content += shape.text + "\n"
            
            return text_content, multimodal_elements
            
        except Exception as e:
            logger.error(f"文档处理失败: {file_path}, 错误: {e}")
            return "", [], f"处理失败: {str(e)}"
    
    def _process_image(self, file_path):
        """处理图像文件"""
        try:
            img = Image.open(file_path)
            return [img]
        except Exception as e:
            logger.warning(f"图像处理失败: {file_path}, 错误: {e}")
            return []
    
    def _process_audio(self, file_path):
        """处理音频文件 - 使用HuggingFace Transformers版本"""
        try:
            if self.model_manager.whisper_model is None or self.model_manager.whisper_processor is None:
                logger.warning("Whisper模型不可用，跳过音频处理")
                return ""
            
            # 加载音频文件
            # import torchaudio
            waveform, sample_rate = torchaudio.load(file_path)
            
            # 重采样到16kHz（Whisper要求的采样率）
            if sample_rate != 16000:
                transform = torchaudio.transforms.Resample(orig_freq=sample_rate, new_freq=16000)
                waveform = transform(waveform)
            
            # 处理音频
            input_features = self.model_manager.whisper_processor(
                waveform.squeeze().numpy(), 
                sampling_rate=16000, 
                return_tensors="pt"
            ).input_features
            
            # 生成转录
            with torch.no_grad():
                predicted_ids = self.model_manager.whisper_model.generate(input_features)
            
            # 解码转录结果
            transcription = self.model_manager.whisper_processor.batch_decode(
                predicted_ids, 
                skip_special_tokens=True
            )[0]
            
            return transcription
            
        except Exception as e:
            logger.error(f"音频处理失败: {file_path}, 错误: {e}")
            return ""
    
    def _process_video(self, file_path):
        """处理视频文件"""
        temp_files = []
        try:
            clip = VideoFileClip(file_path)
            
            # 提取视频帧
            duration = clip.duration
            frame_times = [t for t in range(0, int(duration), 5)][:10]
            frames = []
            for t in frame_times:
                try:
                    frame = clip.get_frame(t)
                    frames.append(frame)
                except:
                    continue
            
            # 提取音频
            audio_path = tempfile.NamedTemporaryFile(suffix='.mp3', delete=False).name
            temp_files.append(audio_path)
            clip.audio.write_audiofile(audio_path, verbose=False, logger=None)
            
            # 转录音频
            audio_text = self._process_audio(audio_path)
            
            # 转换帧为PIL图像
            multimodal_elements = [Image.fromarray(frame) for frame in frames]
            
            return audio_text, multimodal_elements, temp_files
            
        except Exception as e:
            logger.error(f"视频处理失败: {file_path}, 错误: {e}")
            return "", [], temp_files
        finally:
            if 'clip' in locals():
                clip.close()
    
    def _process_text(self, file_path):
        """处理文本文件"""
        try:
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            return ""
        except Exception as e:
            logger.error(f"文本文件处理失败: {file_path}, 错误: {e}")
            return ""

class VectorizationProcessor:
    """向量化处理器"""
    
    def __init__(self):
        self.config_manager = ConfigManager()
        self.model_manager = ModelManager()
        self.file_processor = FileProcessor(self.config_manager)
        self.quality_validator = VectorQualityValidator()
        
        # 确保模型已加载 - 增强错误处理
        self._ensure_models_loaded()
    
    def _ensure_models_loaded(self):
        """确保模型已正确加载"""
        try:
            health = self.model_manager.health_check()
            logger.info(f"模型健康状态: {health}")
            
            # 检查关键模型是否加载
            critical_models_loaded = health.get('bert_model', False) and health.get('bert_tokenizer', False)
            
            if not critical_models_loaded:
                logger.warning("关键模型未加载，重新加载模型...")
                self.model_manager.load_models()
            else:
                logger.info("模型已正确加载")
                
        except Exception as e:
            logger.error(f"模型健康检查失败: {e}，尝试重新加载...")
            try:
                self.model_manager.load_models()
            except Exception as load_error:
                logger.error(f"模型重新加载失败: {load_error}")
                raise
            
    def process_files(self, file_paths):
        """处理多个文件"""
        if not isinstance(file_paths, list):
            file_paths = [file_paths]
        
        # 验证文件
        valid_files = []
        for file_path in file_paths:
            valid, message = self.file_processor.validate_file(file_path)
            if valid:
                valid_files.append(file_path)
            else:
                logger.warning(f"文件验证失败: {file_path} - {message}")
        
        if not valid_files:
            raise ValueError("没有有效的文件可处理")
        
        # 提取内容
        combined_text = ""
        combined_elements = []
        
        for file_path in valid_files:
            text_content, elements, message = self.file_processor.extract_content(file_path)
            combined_text += text_content + "\n"
            combined_elements.extend(elements)
        
        # 生成向量
        vector = self.generate_multimodal_vector(combined_text.strip(), combined_elements)
        
        # 验证向量质量
        is_valid, quality_message = self.quality_validator.validate_vector(vector, combined_text)
        if not is_valid:
            logger.warning(f"向量质量验证失败: {quality_message}")
            # 可以尝试回退到纯文本向量
            vector = self.generate_text_vector(combined_text)
        
        return {
            'file_paths': valid_files,
            'text_content': combined_text,
            'text_length': len(combined_text),
            'element_count': len(combined_elements),
            'vector': vector,
            'vector_dimension': len(vector),
            'quality_status': quality_message if not is_valid else "合格"
        }
    
    def generate_multimodal_vector(self, text_content, multimodal_elements=[]):
        """生成多模态向量"""
        try:
            # 文本向量
            text_embedding = self._generate_text_embedding(text_content)
            
            # 多模态元素向量
            if multimodal_elements and self.model_manager.clip_model:
                multi_embeddings = []
                for elem in multimodal_elements:
                    if isinstance(elem, Image.Image):
                        embedding = self.model_manager.clip_model.encode(elem)
                        multi_embeddings.append(embedding)
                
                if multi_embeddings:
                    avg_multi_emb = np.mean(multi_embeddings, axis=0)
                    norm = np.linalg.norm(avg_multi_emb)
                    if norm > 0:
                        avg_multi_emb = avg_multi_emb / norm
                    
                    # 动态权重调整
                    weight = self.config_manager.get('text_weight')
                    element_count = len(multimodal_elements)
                    if element_count > 5:
                        adjust = min(self.config_manager.get('multi_element_adjust') * (element_count - 5), 0.3)
                        weight = max(weight - adjust, 0.3)
                    
                    final_vector = weight * text_embedding + (1 - weight) * avg_multi_emb
                    norm_final = np.linalg.norm(final_vector)
                    if norm_final > 0:
                        final_vector = final_vector / norm_final
                    else:
                        final_vector = text_embedding
                else:
                    final_vector = text_embedding
            else:
                final_vector = text_embedding
            
            return final_vector.tolist()
            
        except Exception as e:
            logger.error(f"多模态向量生成失败: {e}")
            # 回退到纯文本向量
            return self.generate_text_vector(text_content)
    
    def _generate_text_embedding(self, text_content):
        """生成文本嵌入"""
        if not text_content or not text_content.strip():
            return np.zeros(768)
        
        encoded_input = self.model_manager.bert_tokenizer(
            text_content, 
            padding=True, 
            truncation=True, 
            max_length=512, 
            return_tensors='pt'
        )
        with torch.no_grad():
            model_output = self.model_manager.bert_model(**encoded_input)
        text_embedding = self._mean_pooling(model_output, encoded_input['attention_mask'])
        text_embedding = text_embedding / text_embedding.norm(p=2, dim=-1, keepdim=True)
        return text_embedding.cpu().numpy()[0]
    
    def generate_text_vector(self, text_content):
        """生成纯文本向量"""
        try:
            if not text_content or not text_content.strip():
                return np.zeros(768).tolist()
            
            segments = self._segment_text(text_content)
            segment_vectors = []
            
            for segment in segments:
                encoded_input = self.model_manager.bert_tokenizer(
                    segment, 
                    padding=True, 
                    truncation=True, 
                    max_length=512, 
                    return_tensors='pt'
                )
                with torch.no_grad():
                    model_output = self.model_manager.bert_model(**encoded_input)
                vec = self._mean_pooling(model_output, encoded_input['attention_mask'])
                vec = vec / vec.norm(p=2, dim=-1, keepdim=True)
                segment_vectors.append(vec.cpu().numpy()[0])
            
            if segment_vectors:
                final_vector = np.mean(segment_vectors, axis=0)
                norm = np.linalg.norm(final_vector)
                if norm > 0:
                    final_vector = final_vector / norm
                return final_vector.tolist()
            else:
                return np.zeros(768).tolist()
                
        except Exception as e:
            logger.error(f"文本向量生成失败: {e}")
            return np.zeros(768).tolist()
    
    def _segment_text(self, text):
        """中文文本分词和智能截断"""
        if not text or not text.strip():
            return []
        
        max_length = self.config_manager.get('max_segment_length')
        segments = []
        lines = text.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if len(line) <= max_length:
                segments.append(line)
            else:
                sentences = re.split(r'[。！？!?]', line)
                current_segment = ""
                
                for sentence in sentences:
                    sentence = sentence.strip()
                    if not sentence:
                        continue
                        
                    if len(current_segment) + len(sentence) <= max_length:
                        current_segment += sentence + "。"
                    else:
                        if current_segment:
                            segments.append(current_segment)
                        current_segment = sentence + "。"
                
                if current_segment:
                    segments.append(current_segment)
        
        return segments
    
    def _mean_pooling(self, model_output, attention_mask):
        """BERT向量池化"""
        token_embeddings = model_output[0]
        input_mask_expanded = attention_mask.unsqueeze(-1).expand(token_embeddings.size()).float()
        return torch.sum(token_embeddings * input_mask_expanded, 1) / torch.clamp(input_mask_expanded.sum(1), min=1e-9)

# Celery应用配置
app = Celery('vectorize', broker='redis://localhost:6379/0')
app.conf.update(
    task_serializer='json',
    accept_content=['json'],
    result_serializer='json',
    timezone='Asia/Shanghai',
    enable_utc=True,
    task_routes={
        'vectorize.process_file': {'queue': 'file_processing'},
        'vectorize.process_structured_data': {'queue': 'structured_data'},
        'vectorize.health_check': {'queue': 'monitoring'}
    },
    task_soft_time_limit=300,  # 5分钟软超时
    task_time_limit=600,       # 10分钟硬超时
)

class KafkaManager:
    """Kafka管理器"""
    
    @staticmethod
    def create_producer():
        """创建Kafka生产者 - 使用正确的服务器地址"""
        try:
            
            # 使用您的公网IP地址
            return KafkaProducer(
                bootstrap_servers=['39.108.11.214:9092'],  # 修改为您的公网IP
                value_serializer=lambda v: json.dumps(v, ensure_ascii=False).encode('utf-8'),
                retries=2,
                request_timeout_ms=15000,
                api_version=(2, 0, 2)
            )
        except Exception as e:
            logger.error(f"Kafka生产者创建失败: {e}")
            return None
    
    @staticmethod
    def send_message(producer, topic, message):
        """发送消息"""
        if producer is None:
            return False
            
        try:
            future = producer.send(topic, message)
            future.get(timeout=15)
            logger.info(f"消息发送成功: {topic}")
            return True
        except Exception as e:
            logger.error(f"消息发送失败: {e}")
            return False

class RedisManager:
    """增强的Redis管理器"""
    
    def __init__(self):
        self.connection_pool = redis.ConnectionPool(
            host='localhost',
            port=6379,
            decode_responses=True,
            max_connections=20
        )
    
    def get_connection(self):
        return redis.Redis(connection_pool=self.connection_pool)
    
    def update_task_status(self, task_id, status, result=None, error=None):
        """增强的任务状态更新"""
        try:
            redis_client = self.get_connection()
            task_key = f"task:{task_id}"
            update_data = {
                "status": status,
                "update_time": datetime.now().isoformat()
            }
            
            if status == 'processing':
                update_data["start_time"] = datetime.now().isoformat()
            elif status == 'completed' and result:
                update_data["result"] = json.dumps(result)
                update_data["vector_count"] = result.get('vector_dimension', 0)
                update_data["end_time"] = datetime.now().isoformat()
            elif status == 'failed' and error:
                update_data["error"] = error
                update_data["end_time"] = datetime.now().isoformat()
            
            redis_client.hset(task_key, mapping=update_data)
            redis_client.expire(task_key, 86400)
            return True
        except Exception as e:
            logger.error(f"Redis状态更新失败: {e}")
            return False
    
    def get_task_info(self, task_id):
        """获取任务信息"""
        try:
            redis_client = self.get_connection()
            task_key = f"task:{task_id}"
            return redis_client.hgetall(task_key)
        except Exception as e:
            logger.error(f"获取任务信息失败: {e}")
            return {}

# 全局管理器实例
redis_manager = RedisManager()

@app.task(bind=True, name='vectorize.process_file')
def process_file_task(self, *args, **kwargs):
    """处理文件的Celery任务"""
    # 解析参数
    if len(args) >= 2:
        file_paths = args[0]
        task_id = args[1] if len(args) > 1 else None
    else:
        file_paths = kwargs.get('file_paths', [])
        task_id = kwargs.get('task_id', None)
        
    processor = VectorizationProcessor()
    redis_manager = RedisManager()
    
    # 添加模型健康检查
    model_health = processor.model_manager.health_check()
    if not model_health.get('bert_model') or not model_health.get('bert_tokenizer'):
        logger.warning("BERT模型未加载，强制重新加载...")
        processor.model_manager.load_models()
    
    # 更新任务状态为处理中
    redis_manager.update_task_status(task_id, 'processing')
    
    try:
        logger.info(f"开始处理文件任务 {task_id}: {file_paths}")
        
        result = processor.process_files(file_paths)
        result.update({
            'task_id': task_id or self.request.id,
            'status': 'success',
            'timestamp': datetime.now().isoformat(),
            'modality': 'moltimodal'
        })
        
        # 发送到Kafka
        kafka_result = _send_to_kafka(
            result['vector'], 
            result['text_content'], 
            file_paths, 
            task_id
        )
        result['kafka_status'] = 'sent' if kafka_result else 'failed'
        
        # 更新Redis任务状态为完成
        redis_manager.update_task_status(task_id, 'completed', result)
        
        logger.info(f"文件处理完成: {file_paths}, 向量维度: {len(result['vector'])}")
        return result
        
    except Exception as e:
        logger.error(f"文件处理任务失败: {e}")
        error_result = {'error': str(e), 'task_id': task_id}
        redis_manager.update_task_status(task_id, 'failed', error=str(e))
        raise self.retry(exc=e, countdown=60, max_retries=3)

# 添加任务状态查询函数
@app.task(bind=True, name='vectorize.get_task_status')
def get_task_status_task(self, task_id):
    """获取任务状态的Celery任务"""
    redis_manager = RedisManager()
    return redis_manager.get_task_info(task_id)
    
# 添加批量任务状态查询
@app.task(bind=True, name='vectorize.batch_task_status')
def batch_task_status_task(self, task_ids):
    """批量获取任务状态"""
    redis_manager = RedisManager()
    results = {}
    for task_id in task_ids:
        results[task_id] = redis_manager.get_task_info(task_id)
    return results

def _send_to_kafka(vector, text_content, file_paths, task_id=None):
    """发送向量数据到Kafka"""
    try:
        producer = KafkaManager.create_producer()
        if not producer:
            return False
            
        file_id = task_id or hashlib.md5(','.join(file_paths).encode()).hexdigest()
        
        message = {
            "id": file_id,
            "type": "multimodal_document",
            "path": file_paths[0] if file_paths else "unknown",
            "vector": vector,
            "content": text_content[:2048],
            "metadata": {
                "file_count": len(file_paths),
                "content_length": len(text_content),
                "timestamp": datetime.now().isoformat(),
                "vector_dimension": len(vector)
            }
        }
        
        # 添加校验和
        # message['checksum'] = hashlib.md5(
        #     json.dumps(message, sort_keys=True).encode()
        # ).hexdigest()
        
        # 计算校验和
        message_for_checksum = message.copy()
        checksum_data = json.dumps(message_for_checksum, sort_keys=True, ensure_ascii=False)
        message['checksum'] = hashlib.md5(checksum_data.encode('utf-8')).hexdigest()
        
        success = KafkaManager.send_message(producer, 'vector-data', message)
        producer.close()
        
        if success:
            logger.info(f"向量数据发送到Kafka成功: {file_id}")
        else:
            logger.error(f"向量数据发送到Kafka失败: {file_id}")
            
        return success
        
    except Exception as e:
        logger.error(f"Kafka发送失败: {e}")
        return False

@app.task(bind=True, name='vectorize.process_structured_data')
def process_structured_data_task(self, data_type, data_id, meta_data):
    """处理结构化数据的Celery任务"""
    try:
        # 输入验证
        if not data_type or not data_id:
            raise ValueError("数据类型和ID不能为空")
        
        if isinstance(meta_data, str):
            try:
                meta_data = json.loads(meta_data)
            except json.JSONDecodeError:
                raise ValueError("元数据JSON格式错误")
        
        # 检查重复处理
        redis_client = redis_manager.get_connection()
        if redis_client.exists(f"processed:{data_type}:{data_id}"):
            logger.info(f"跳过已处理的数据: {data_type}_{data_id}")
            return {"status": "skipped", "reason": "already_processed"}
        
        processor = VectorizationProcessor()
        
        # 根据数据类型处理
        if data_type == "product":
            result = _process_product_data(processor, data_id, meta_data)
        elif data_type == "user_query":
            result = _process_user_query(processor, data_id, meta_data)
        else:
            raise ValueError(f"不支持的数据类型: {data_type}")
        
        # 标记为已处理
        redis_client.setex(f"processed:{data_type}:{data_id}", 86400, "1")
        
        logger.info(f"结构化数据处理完成: {data_type}_{data_id}")
        return result
        
    except Exception as e:
        logger.error(f"结构化数据处理失败: {e}")
        raise self.retry(exc=e, countdown=30, max_retries=3)

def _process_product_data(processor, product_id, meta_data):
    """处理商品数据"""
    try:
        # 构建商品描述文本
        text_parts = [
            f"商品ID: {product_id}",
            f"商品名称: {meta_data.get('goods_name', '')}",
            f"店铺名称: {meta_data.get('shop_name', '')}",
            f"商品分类: {meta_data.get('category_name', '')}",
            f"商品描述: {meta_data.get('description', '')}"
        ]
        
        # 处理规格信息
        sku_vectors = []
        valid_skus = []
        
        for sku in meta_data.get('all_skus', []):
            sku_text = f"规格: {sku.get('sku_name', '')} 价格: {sku.get('price', 0)} 库存: {sku.get('goods_number', 0)}"
            sku_vector = processor.generate_text_vector(sku_text)
            if sku_vector and len(sku_vector) == 768:
                sku_vectors.append(sku_vector)
                valid_skus.append(sku)
        
        # 生成主商品向量
        if sku_vectors:
            main_vector = np.mean(sku_vectors, axis=0).tolist()
        else:
            combined_text = " ".join(text_parts)
            main_vector = processor.generate_text_vector(combined_text)
        
        # 准备消息
        messages = []
        product_message = {
            "id": f"product_{product_id}",
            "type": "product",
            "vector": main_vector,
            "metadata": {
                "goods_name": meta_data.get('goods_name', ''),
                "shop_name": meta_data.get('shop_name', ''),
                "category": meta_data.get('category_name', ''),
                "sku_count": len(valid_skus),
                "timestamp": datetime.now().isoformat()
            }
        }
        messages.append(product_message)
        
        # 规格消息
        for i, sku in enumerate(valid_skus):
            sku_message = {
                "id": f"sku_{sku.get('sku_id', i)}",
                "type": "sku",
                "vector": sku_vectors[i],
                "parent_id": f"product_{product_id}",
                "metadata": {
                    "sku_name": sku.get('sku_name', ''),
                    "price": sku.get('price', 0),
                    "currency": sku.get('currency', 'CNY'),
                    "stock": sku.get('goods_number', 0)
                }
            }
            messages.append(sku_message)
        
        # 批量发送
        success_count = 0
        producer = KafkaManager.create_producer()
        
        for message in messages:
            if KafkaManager.send_message(producer, 'vector-data', message):
                success_count += 1
        
        producer.close()
        
        return {
            "status": "success",
            "product_id": product_id,
            "sku_count": len(valid_skus),
            "messages_sent": success_count,
            "total_messages": len(messages)
        }
        
    except Exception as e:
        logger.error(f"商品数据处理失败: {e}")
        raise

def _process_user_query(processor, query_id, meta_data):
    """处理用户查询"""
    try:
        query_text = meta_data.get('query', '')
        if not query_text:
            raise ValueError("查询文本不能为空")
        
        query_vector = processor.generate_text_vector(query_text)
        
        message = {
            "id": f"query_{query_id}",
            "type": "user_query",
            "vector": query_vector,
            "content": query_text,
            "metadata": {
                "user_id": meta_data.get('user_id', ''),
                "session_id": meta_data.get('session_id', ''),
                "timestamp": datetime.now().isoformat(),
                "query_length": len(query_text)
            }
        }
        
        producer = KafkaManager.create_producer()
        success = KafkaManager.send_message(producer, 'vector-data', message)
        producer.close()
        
        return {
            "status": "success" if success else "failed",
            "query_id": query_id,
            "vector_generated": True
        }
            
    except Exception as e:
        logger.error(f"用户查询处理失败: {e}")
        raise

@app.task(bind=True, name='vectorize.health_check')
def health_check_task(self):
    """系统健康检查任务"""
    try:
        model_manager = ModelManager()
        check_results = {}
        
        # 检查模型
        check_results['models'] = model_manager.health_check()
        
        # 检查Redis
        try:
            redis_client = redis_manager.get_connection()
            redis_client.ping()
            check_results['redis'] = {'status': 'healthy'}
        except Exception as e:
            check_results['redis'] = {'status': 'unhealthy', 'error': str(e)}
        
        # 检查Kafka
        try:
            producer = KafkaManager.create_producer()
            if producer:
                producer.close()
                check_results['kafka'] = {'status': 'healthy'}
            else:
                check_results['kafka'] = {'status': 'unhealthy', 'error': 'Producer creation failed'}
        except Exception as e:
            check_results['kafka'] = {'status': 'unhealthy', 'error': str(e)}
        
        # 检查系统资源
        try:
            disk_usage = psutil.disk_usage('/')
            memory_usage = psutil.virtual_memory()
            cpu_usage = psutil.cpu_percent(interval=1)
            
            check_results['system'] = {
                'disk': {
                    'total_gb': round(disk_usage.total / (1024**3), 2),
                    'used_gb': round(disk_usage.used / (1024**3), 2),
                    'free_gb': round(disk_usage.free / (1024**3), 2),
                    'usage_percent': disk_usage.percent
                },
                'memory': {
                    'usage_percent': memory_usage.percent,
                    'available_gb': round(memory_usage.available / (1024**3), 2)
                },
                'cpu': {
                    'usage_percent': cpu_usage
                }
            }
        except Exception as e:
            check_results['system'] = {'error': str(e)}
        
        # 总体状态
        all_healthy = all([
            all(check_results['models'].values()),
            check_results['redis']['status'] == 'healthy',
            check_results['kafka']['status'] == 'healthy',
            check_results['system'].get('disk', {}).get('usage_percent', 100) < 90,
            check_results['system'].get('memory', {}).get('usage_percent', 100) < 90
        ])
        
        check_results['overall_status'] = 'healthy' if all_healthy else 'degraded'
        check_results['timestamp'] = datetime.now().isoformat()
        
        return check_results
        
    except Exception as e:
        logger.error(f"健康检查失败: {e}")
        return {'overall_status': 'error', 'error': str(e)}

class MonitoringService:
    """监控服务"""
    
    def __init__(self):
        self.redis_manager = redis_manager
    
    def get_task_statistics(self, hours=24):
        """获取任务统计"""
        try:
            redis_client = self.redis_manager.get_connection()
            task_keys = redis_client.keys("task:*")
            
            cutoff_time = datetime.now() - timedelta(hours=hours)
            recent_tasks = []
            
            for key in task_keys:
                task_data = redis_client.hgetall(key)
                create_time_str = task_data.get('create_time', '')
                
                if create_time_str:
                    try:
                        create_time = datetime.fromisoformat(create_time_str)
                        if create_time > cutoff_time:
                            recent_tasks.append(task_data)
                    except:
                        continue
            
            stats = {
                'total_tasks': len(recent_tasks),
                'completed': len([t for t in recent_tasks if t.get('status') == 'completed']),
                'failed': len([t for t in recent_tasks if t.get('status') == 'failed']),
                'pending': len([t for t in recent_tasks if t.get('status') == 'pending']),
                'processing': len([t for t in recent_tasks if t.get('status') == 'processing'])
            }
            
            if stats['total_tasks'] > 0:
                stats['success_rate'] = round(stats['completed'] / stats['total_tasks'] * 100, 2)
            else:
                stats['success_rate'] = 0
                
            return stats
            
        except Exception as e:
            logger.error(f"获取任务统计失败: {e}")
            return {}
    
    def get_performance_metrics(self):
        """获取性能指标"""
        try:
            return {
                'vector_dimension': 768,
                'average_processing_time': '待实现',
                'throughput': '待实现'
            }
        except Exception as e:
            logger.error(f"获取性能指标失败: {e}")
            return {}

def process_files_sync(file_paths):
    processor = VectorizationProcessor()
    try:
        result = processor.process_files(file_paths)
        
        # 修复：直接检查处理结果，不依赖'success'键
        if result and isinstance(result, dict):
            # 如果process_files返回了有效结果，直接返回
            if 'vector' in result and result['vector']:
                return {
                    'success': True,
                    'vector': result['vector'],
                    'text_content': result.get('text_content', ''),
                    'text_length': result.get('text_length', 0)
                }
            else:
                error_msg = result.get('error', '向量生成失败')
                return {'success': False, 'error': error_msg}
        else:
            return {'success': False, 'error': '处理结果格式错误'}
            
    except Exception as e:
        logger.error(f"处理失败: {e}")
        return {'success': False, 'error': str(e)}

def mark_task_completed_in_sync_mode(task_id, result):
    try:
        r = redis.Redis(host='127.0.0.1', port=6379, db=0, decode_responses=True)
        r.hset(f"task:{task_id}", mapping={
            "status": "completed",
            "vector_dimension": len(result.get('vector', [])),
            "text_length": result.get('text_length', 0),
            "process_time": datetime.now().isoformat()
        })
        r.expire(f"task:{task_id}", 3600)
        logger.info(f"Sync mode: task {task_id} marked as completed in Redis")  # 改用日志器
    except Exception as e:
        logger.error(f"Redis update failed in sync mode: {e}")  # 改用日志器

# 修复Jieba缓存权限问题
def fix_jieba_permission():
    try:
        # 设置Jieba缓存到用户目录
        jieba_cache_dir = os.path.expanduser('~/.jieba_cache')
        os.makedirs(jieba_cache_dir, exist_ok=True)
        jieba_cache_file = os.path.join(jieba_cache_dir, 'jieba.cache')
        
        # 设置环境变量
        os.environ['JIEBA_CACHE'] = jieba_cache_file
        
        # 或者在代码中直接设置
        import jieba
        jieba.dt.cache_file = jieba_cache_file
        
    except Exception as e:
        logger.warning(f"Jieba缓存设置失败: {e}")

# 主程序入口
if __name__ == "__main__":
    # 初始化jieba（静默模式）
    fix_jieba_permission()
    
    try:
        # 设置jieba为静默模式
        jieba.setLogLevel(logging.ERROR)
    except:
        pass
    
    if len(sys.argv) > 1 and sys.argv[1] == "sync_process":
        
        # 重定向stdout到/dev/null
        original_stdout = sys.stdout
        sys.stdout = open(os.devnull, 'w')
        
        try:
            if len(sys.argv) > 2:
                input_data = json.loads(sys.argv[2])
                file_path = input_data[0]
                task_id = input_data[1]
                result = process_files_sync([file_path])
                
                # 恢复stdout并输出纯JSON
                sys.stdout = original_stdout
                print(json.dumps(result, ensure_ascii=False))
                
                mark_task_completed_in_sync_mode(task_id, result)
            else:
                sys.stdout = original_stdout
                print(json.dumps({'success': False, 'error': '参数错误'}))
        except Exception as e:
            sys.stdout = original_stdout
            print(json.dumps({'success': False, 'error': str(e)}))
        finally:
            # 确保stdout被恢复
            if 'original_stdout' in locals():
                sys.stdout = original_stdout
        sys.exit(0)
    else:
        if len(sys.argv) > 1:
            command = sys.argv[1]
            
            if command == "health_check":
                result = health_check_task.apply()
                print(json.dumps(result.result, indent=2))
                
            elif command == "stats":
                monitor = MonitoringService()
                stats = monitor.get_task_statistics()
                print(json.dumps(stats, indent=2))
                
            elif command == "vectorize_file":
                if len(sys.argv) > 2:
                    file_paths = json.loads(sys.argv[2])
                    result = process_file_task.apply(args=[file_paths])
                    print(json.dumps(result.result, indent=2))
                else:
                    print("Usage: python vectorize.py vectorize_file '[\"path1\", \"path2\"]'")
                    
            elif command == "worker":
                try:
                    app.start(argv=['worker', '--loglevel=info', '--pool=prefork', '--concurrency=2'])
                except KeyboardInterrupt:
                    logger.info("Worker停止")
                except Exception as e:
                    logger.error(f"Worker启动失败: {e}")
                    sys.exit(1)
            else:
                print("未知命令")
                print("可用命令: health_check, stats, vectorize_file, worker")
        else:
            try:
                app.start(argv=['worker', '--loglevel=info', '--pool=prefork', '--concurrency=2'])
            except Exception as e:
                logger.error(f"Worker启动失败: {e}")
                sys.exit(1)